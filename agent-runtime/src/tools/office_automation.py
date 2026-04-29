"""
Office 文档自动化工具
提供 Excel、Word、PowerPoint 文档生成能力
"""
import json
import os
import tempfile
from typing import Dict, Any, Optional
from pathlib import Path

try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
    
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt
    from docx.oxml.ns import qn
    
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    import pandas as pd
    
    OFFICE_AVAILABLE = True
except ImportError:
    OFFICE_AVAILABLE = False


class OfficeAutomationTool:
    """Office 文档自动化工具类"""
    
    def __init__(self):
        if not OFFICE_AVAILABLE:
            raise ImportError(
                "Office automation libraries not available. "
                "Please install with: pip install openpyxl python-docx python-pptx pandas"
            )
    
    def generate_excel(self, data: Dict[str, Any], output_path: str, title: str = "报表") -> str:
        """
        生成 Excel 报表
        
        Args:
            data: 包含数据的字典
            output_path: 输出文件路径
            title: 报表标题
            
        Returns:
            生成的文件路径
        """
        wb = Workbook()
        ws = wb.active
        ws.title = "数据报表"
        
        # 设置标题
        ws['A1'] = title
        ws['A1'].font = Font(size=16, bold=True)
        ws.merge_cells('A1:D1')
        
        row = 3
        
        # 处理数据集
        if 'datasets' in data:
            for dataset_name, dataset_data in data['datasets'].items():
                # 写入数据集标题
                ws.cell(row=row, column=1, value=dataset_name)
                ws.cell(row=row, column=1).font = Font(bold=True)
                row += 1
                
                if isinstance(dataset_data, list) and dataset_data:
                    if isinstance(dataset_data[0], dict):
                        # 如果是字典列表，转换为 DataFrame
                        df = pd.DataFrame(dataset_data)
                        
                        # 写入列标题
                        headers = df.columns.tolist()
                        for col_idx, header in enumerate(headers, 1):
                            cell = ws.cell(row=row, column=col_idx, value=header)
                            cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
                            cell.font = Font(color="FFFFFF", bold=True)
                            cell.alignment = Alignment(horizontal="center")
                        
                        row += 1
                        
                        # 写入数据行
                        for _, data_row in df.iterrows():
                            for col_idx, value in enumerate(data_row.values, 1):
                                ws.cell(row=row, column=col_idx, value=value)
                            row += 1
                    else:
                        # 简单列表数据
                        for data_row in dataset_data:
                            for col_idx, value in enumerate(data_row, 1):
                                ws.cell(row=row, column=col_idx, value=value)
                            row += 1
                
                row += 2  # 空行分隔
        
        # 添加边框
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        for row in ws.iter_rows(min_row=2, max_col=ws.max_column, max_row=ws.max_row):
            for cell in row:
                if cell.value is not None:
                    cell.border = thin_border
        
        wb.save(output_path)
        return output_path
    
    def generate_word(self, data: Dict[str, Any], output_path: str) -> str:
        """
        生成 Word 文档
        
        Args:
            data: 包含文档数据的字典
            output_path: 输出文件路径
            
        Returns:
            生成的文件路径
        """
        doc = Document()
        
        # 添加标题
        if 'title' in data:
            title_para = doc.add_heading(data['title'], 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 添加副标题
        if 'subtitle' in data:
            subtitle_para = doc.add_paragraph()
            subtitle_para.add_run(data['subtitle']).italic = True
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 添加段落
        if 'sections' in data:
            for section in data['sections']:
                if 'title' in section:
                    heading_level = section.get('heading_level', 1)
                    doc.add_heading(section['title'], level=heading_level)
                
                if 'content' in section:
                    if isinstance(section['content'], list):
                        for para_text in section['content']:
                            para = doc.add_paragraph()
                            para.add_run(str(para_text))
                    else:
                        para = doc.add_paragraph()
                        para.add_run(str(section['content']))
        
        # 添加表格
        if 'tables' in data:
            for table_data in data['tables']:
                if 'rows' in table_data and table_data['rows']:
                    # 创建表格
                    rows = len(table_data['rows'])
                    cols = len(table_data['rows'][0]) if table_data['rows'] else 0
                    table = doc.add_table(rows=rows, cols=cols)
                    table.style = 'Table Grid'
                    
                    # 填充数据
                    for i, row_data in enumerate(table_data['rows']):
                        for j, cell_data in enumerate(row_data):
                            cell = table.cell(i, j)
                            cell.text = str(cell_data)
                            
                            # 如果是标题行，加粗
                            if i == 0:
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
        
        doc.save(output_path)
        return output_path
    
    def generate_powerpoint(self, data: Dict[str, Any], output_path: str) -> str:
        """
        生成 PowerPoint 演示文稿
        
        Args:
            data: 包含演示文稿数据的字典
            output_path: 输出文件路径
            
        Returns:
            生成的文件路径
        """
        prs = Presentation()
        
        # 设置幻灯片大小
        prs.slide_width = PptxInches(13.33)  # 16:9 宽屏
        prs.slide_height = PptxInches(7.5)
        
        # 添加标题幻灯片
        if 'title_slide' in data:
            title_slide_data = data['title_slide']
            title_slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = title_slide_data.get('title', '演示文稿标题')
            subtitle.text = title_slide_data.get('subtitle', '副标题')
        
        # 添加内容幻灯片
        if 'slides' in data:
            for slide_data in data['slides']:
                if slide_data['type'] == 'title_and_content':
                    layout = prs.slide_layouts[1]  # 标题和内容布局
                    slide = prs.slides.add_slide(layout)
                    
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_data.get('title', '幻灯片标题')
                    content.text = slide_data.get('content', '幻灯片内容')
                    
                    # 如果有列表项，添加到内容中
                    if 'bullet_points' in slide_data:
                        for point in slide_data['bullet_points']:
                            content.text += f"\n• {point}"
        
        prs.save(output_path)
        return output_path


# 全局实例
office_tool = None


def get_office_tool() -> OfficeAutomationTool:
    """获取 Office 自动化工具实例"""
    global office_tool
    if office_tool is None:
        office_tool = OfficeAutomationTool()
    return office_tool


def create_office_document(document_type: str, data: Dict[str, Any], output_filename: Optional[str] = None) -> Dict[str, Any]:
    """
    创建 Office 文档的统一接口
    
    Args:
        document_type: 文档类型 ('excel', 'word', 'powerpoint')
        data: 文档数据
        output_filename: 输出文件名（可选）
        
    Returns:
        包含生成结果的字典
    """
    tool = get_office_tool()
    
    # 创建临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{document_type}") as tmp_file:
        output_path = tmp_file.name
    
    try:
        if document_type.lower() == 'excel':
            if output_filename is None:
                output_path = output_path.replace('.excel', '.xlsx')
            else:
                output_path = os.path.join(os.path.dirname(output_path), output_filename)
            result_path = tool.generate_excel(data, output_path, data.get('title', '报表'))
        elif document_type.lower() == 'word':
            if output_filename is None:
                output_path = output_path.replace('.word', '.docx')
            else:
                output_path = os.path.join(os.path.dirname(output_path), output_filename)
            result_path = tool.generate_word(data, output_path)
        elif document_type.lower() == 'powerpoint':
            if output_filename is None:
                output_path = output_path.replace('.powerpoint', '.pptx')
            else:
                output_path = os.path.join(os.path.dirname(output_path), output_filename)
            result_path = tool.generate_powerpoint(data, output_path)
        else:
            raise ValueError(f"Unsupported document type: {document_type}")
        
        return {
            "success": True,
            "file_path": result_path,
            "file_size": os.path.getsize(result_path),
            "message": f"{document_type.upper()} 文档已生成: {result_path}"
        }
    
    except Exception as e:
        # 清理临时文件
        if os.path.exists(output_path):
            os.unlink(output_path)
        return {
            "success": False,
            "error": str(e),
            "message": f"生成 {document_type} 文档时出错: {str(e)}"
        }


if __name__ == "__main__":
    # 测试代码
    test_data = {
        "datasets": {
            "销售数据": [
                {"月份": "1月", "销售额": 10000, "成本": 7000, "利润": 3000},
                {"月份": "2月", "销售额": 12000, "成本": 8000, "利润": 4000},
                {"月份": "3月", "销售额": 15000, "成本": 9000, "利润": 6000}
            ]
        }
    }
    
    result = create_office_document('excel', test_data, 'test_sales_report.xlsx')
    print(result)